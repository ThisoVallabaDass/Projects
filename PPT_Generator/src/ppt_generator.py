"""
================================================================================
KAARTECH CORPORATE PPT GENERATOR — v2.0 (Complete Rewrite)
================================================================================
Author : Thiso Vallabadass K (AID-902015)
Purpose: AI-powered presentation generator that:
         1. Reads the KaarTech Corporate Template layouts (183 layouts)
         2. Uses Gemini 2.0 Flash to generate structured JSON content
         3. Adds slides using prs.slides.add_slide(layout) — NOT XML cloning
         4. Injects content with explicit font overrides at XML run level
         5. Outputs a fully branded KaarTech .pptx file

Usage:
    python ppt_generator.py --topic "SAP S/4HANA Migration" --output "migration.pptx"
    python ppt_generator.py --prompt-file my_prompt.txt --output "my_presentation.pptx"
    python ppt_generator.py --json-file content.json --output "my_presentation.pptx"

Dependencies:
    pip install python-pptx google-generativeai python-dotenv lxml Pillow requests
================================================================================
"""

import os
import sys
import json
import copy
import argparse
import shutil
import tempfile
from pathlib import Path
from typing import Optional, List
from dotenv import load_dotenv

# ── python-pptx ──────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import lxml.etree as etree

# ── Gemini AI ─────────────────────────────────────────────────────────────────
import google.generativeai as genai

# ── Load .env (GEMINI_API_KEY lives here) ─────────────────────────────────────
_env_path = Path(__file__).parent.parent / ".env"
if _env_path.exists():
    load_dotenv(_env_path)
else:
    # Try .env.example as fallback hint
    load_dotenv(Path(__file__).parent.parent / ".env.example")


# ═════════════════════════════════════════════════════════════════════════════
#  CONSTANTS & PATHS
# ═════════════════════════════════════════════════════════════════════════════

BASE_DIR = Path(__file__).parent.parent
TEMPLATE_PATH = BASE_DIR / "reference" / "KaarTech Corporate PPT Templates - 2024 Guide (2).pptx"
OUTPUT_DIR = BASE_DIR / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# ── KaarTech Brand Colors ────────────────────────────────────────────────────
KAAR_RED   = RGBColor(0xC0, 0x20, 0x2B)   # Primary red  #C0202B
KAAR_DARK  = RGBColor(0x1A, 0x1A, 0x2E)   # Dark navy    #1A1A2E
KAAR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)    # White        #FFFFFF
KAAR_GRAY  = RGBColor(0x58, 0x58, 0x58)   # Body gray    #585858
KAAR_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)   # Light bg     #F5F5F5


# ═════════════════════════════════════════════════════════════════════════════
#  LAYOUT MAP — slide_type → layout index in the 183-layout template
# ═════════════════════════════════════════════════════════════════════════════

LAYOUT_MAP = {
    # ── Title / Cover slides ─────────────────────────────────────────────────
    "cover":            0,      # Title Slide _2 — red geometric cover
    "cover_dark":       1,      # Title Slide _3 — dark navy
    "cover_red":        8,      # Title Slide _10 — full dark red
    # ── Agenda slides ────────────────────────────────────────────────────────
    "agenda_6":         30,     # Agenda Slide_1 — 6 items, 2 cols (3+3)
    "agenda_8":         31,     # Agenda Slide_2 — 8 items, pentagon arrows
    "agenda":           31,     # Default agenda → 8-item layout
    # ── Section dividers ─────────────────────────────────────────────────────
    "divider":          20,     # Slide Splitter_1 — dark red full bleed
    "divider_2":        21,     # Slide Splitter_2 — geometric shapes
    "divider_3":        22,     # Slide Splitter_3 — diagonal line accent
    # ── Content slides ───────────────────────────────────────────────────────
    "overview":         14,     # Content Slide_3 — title + 1 tall content box
    "content_single":   14,     # Content Slide_3 — single content area
    "content_two":      15,     # Content Slide_4 — two side-by-side columns
    "content_three":    17,     # Content Slide_6 — three columns
    "comparison":       16,     # Content Slide_5 — two header+content pairs
    # ── Process / Flow slides ────────────────────────────────────────────────
    "process_arrow":    106,    # Arrow Slide_1 — 4-step arrow flow
    "process_flow":     135,    # Flow Diagram_1 — 5-step box flow
    "process_cycle":    84,     # Cycle Slide_1 — circular process
    # ── Timeline ─────────────────────────────────────────────────────────────
    "timeline":         59,     # Timeline Slide_1 — horizontal, 9 steps
    # ── Components ───────────────────────────────────────────────────────────
    "components":       15,     # Content Slide_4 — two columns
    "components_3":     17,     # Content Slide_6 — three columns
    # ── Dashboard ────────────────────────────────────────────────────────────
    "dashboard":        14,     # Content Slide_3 — single column dashboard
    # ── Special layouts ──────────────────────────────────────────────────────
    "pointer_4":        96,     # Pointers Slide_1 — 4 numbered items
    "circle_4":         74,     # Circle Slide_1 — 4-item circle
    "quote":            45,     # Quote Slide_1 — highlight statement
    "table":            116,    # Table Slide_1 — data grid
    # ── Thank You ────────────────────────────────────────────────────────────
    "thankyou":         174,    # Thank You Slide_6 — KaarTech address card
    "thankyou_simple":  169,    # Thank You Slide_1 — abstract shapes
}


# ═════════════════════════════════════════════════════════════════════════════
#  FONT INJECTION UTILITIES
#  These set font at the XML run level to guarantee bold/color/size appear
#  visually, overriding any theme/layout defaults.
# ═════════════════════════════════════════════════════════════════════════════

def _color_val(color_rgb: RGBColor) -> str:
    """Convert RGBColor to 6-char hex string."""
    return f"{color_rgb.red:02X}{color_rgb.green:02X}{color_rgb.blue:02X}"


def force_font(shape, text, font_name="Arial", size_pt=14,
               bold=False, color_rgb=None):
    """
    Set text AND explicitly override font at paragraph + run level.
    Clears ALL existing content and replaces with a single paragraph/run.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True

    # Clear all existing paragraphs except first
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    para = tf.paragraphs[0]

    # Clear existing runs
    for r in para.runs:
        r._r.getparent().remove(r._r)

    # Also clear any direct text elements in the paragraph
    for t_elem in list(para._p.findall(qn('a:r'))):
        para._p.remove(t_elem)

    # Ensure paragraph properties exist
    pPr = para._p.get_or_add_pPr()

    # Create run element with explicit formatting
    r_elem = etree.SubElement(para._p, qn('a:r'))
    rPr = etree.SubElement(r_elem, qn('a:rPr'), attrib={
        'lang': 'en-US', 'dirty': '0'
    })
    rPr.set('b', '1' if bold else '0')
    rPr.set('sz', str(int(size_pt * 100)))

    # Font name
    etree.SubElement(rPr, qn('a:latin'), attrib={'typeface': font_name})

    # Color
    if color_rgb:
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        etree.SubElement(solidFill, qn('a:srgbClr'),
                         attrib={'val': _color_val(color_rgb)})

    # Text content
    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text


def append_run(para, text, font_name="Calibri", size_pt=12,
               bold=False, italic=False, color_rgb=None):
    """Append a run to an existing paragraph with explicit formatting."""
    r_elem = etree.SubElement(para._p, qn('a:r'))
    rPr = etree.SubElement(r_elem, qn('a:rPr'), attrib={
        'lang': 'en-US', 'dirty': '0'
    })
    rPr.set('b', '1' if bold else '0')
    rPr.set('i', '1' if italic else '0')
    rPr.set('sz', str(int(size_pt * 100)))

    etree.SubElement(rPr, qn('a:latin'), attrib={'typeface': font_name})

    if color_rgb:
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        etree.SubElement(solidFill, qn('a:srgbClr'),
                         attrib={'val': _color_val(color_rgb)})

    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text


def add_formatted_paragraph(text_frame, text, font_name="Calibri", size_pt=12,
                            bold=False, italic=False, color_rgb=None, level=0):
    """Add a new paragraph to an existing text frame with explicit formatting."""
    p_elem = etree.SubElement(text_frame._txBody, qn('a:p'))
    pPr = etree.SubElement(p_elem, qn('a:pPr'))
    if level > 0:
        pPr.set('lvl', str(level))

    r_elem = etree.SubElement(p_elem, qn('a:r'))
    rPr = etree.SubElement(r_elem, qn('a:rPr'), attrib={
        'lang': 'en-US', 'dirty': '0'
    })
    rPr.set('b', '1' if bold else '0')
    rPr.set('i', '1' if italic else '0')
    rPr.set('sz', str(int(size_pt * 100)))

    etree.SubElement(rPr, qn('a:latin'), attrib={'typeface': font_name})

    if color_rgb:
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        etree.SubElement(solidFill, qn('a:srgbClr'),
                         attrib={'val': _color_val(color_rgb)})

    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text


def clear_text_body(text_frame):
    """Remove all <a:p> elements from a text frame's txBody."""
    txBody = text_frame._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def add_slide_from_layout(prs, layout_idx):
    """Add a new slide using the specified layout index."""
    try:
        layout = prs.slide_layouts[layout_idx]
    except IndexError:
        print(f"[WARN] Layout index {layout_idx} out of range "
              f"(max {len(prs.slide_layouts)-1}), falling back to layout 14.")
        layout = prs.slide_layouts[14]
    return prs.slides.add_slide(layout)


def get_ph(slide, ph_idx):
    """Get a placeholder shape by its placeholder format idx."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            return shape
    return None


def get_text_shapes_by_area(slide, exclude_shape=None):
    """Return text shapes sorted by area (largest first), excluding one."""
    shapes = []
    for s in slide.shapes:
        if s.has_text_frame:
            if exclude_shape is not None and s.shape_id == exclude_shape.shape_id:
                continue
            shapes.append(s)
    shapes.sort(key=lambda s: s.width * s.height, reverse=True)
    return shapes


# ═════════════════════════════════════════════════════════════════════════════
#  SMART LAYOUT SELECTION
# ═════════════════════════════════════════════════════════════════════════════

def choose_layout(slide_data, agenda=None):
    """
    Choose the correct layout index based on slide_data.
    Rules:
      1. layout_override → use directly
      2. slide_type → LAYOUT_MAP
      3. process sub-selection by step count / is_cycle
      4. agenda sub-selection by item count
      5. components sub-selection by count
      6. Default fallback → layout 14
    """
    # Rule 1: explicit override
    if "layout_override" in slide_data and slide_data["layout_override"] is not None:
        return int(slide_data["layout_override"])

    stype = slide_data.get("slide_type", "overview")

    # Rule 3: process sub-selection
    if stype == "process":
        steps = slide_data.get("process_steps", [])
        if slide_data.get("is_cycle"):
            return LAYOUT_MAP["process_cycle"]   # layout 84
        elif len(steps) <= 4:
            return LAYOUT_MAP["process_arrow"]   # layout 106
        elif len(steps) == 5:
            return LAYOUT_MAP["process_flow"]    # layout 135
        else:
            return LAYOUT_MAP["process_flow"]    # layout 135

    # Rule 4: agenda sub-selection
    if stype == "agenda":
        items = slide_data.get("agenda", agenda or [])
        if len(items) <= 6:
            return LAYOUT_MAP["agenda_6"]   # layout 30
        else:
            return LAYOUT_MAP["agenda_8"]   # layout 31

    # Rule 5: components sub-selection
    if stype == "components":
        comps = slide_data.get("components", [])
        if len(comps) <= 2:
            return LAYOUT_MAP["content_single"]   # layout 14
        elif len(comps) == 3:
            return LAYOUT_MAP["components_3"]     # layout 17
        else:
            return LAYOUT_MAP["components"]       # layout 15

    # Rule 2: direct LAYOUT_MAP lookup
    if stype in LAYOUT_MAP:
        return LAYOUT_MAP[stype]

    # Rule 6: default fallback
    return 14


# ═════════════════════════════════════════════════════════════════════════════
#  PER-LAYOUT FILL FUNCTIONS
# ═════════════════════════════════════════════════════════════════════════════

# ── COVER SLIDE ──────────────────────────────────────────────────────────────

def fill_cover(slide, data):
    """
    Fill cover slide. Layout 0 has idx=0 "Title 1" placeholder.
    Find second-largest text shape for subtitle/presenter line.
    """
    title_text = data.get("title", "KaarTech Presentation")

    # Title placeholder: try idx=0, then idx=11
    ph_title = get_ph(slide, 0) or get_ph(slide, 11)
    if ph_title:
        force_font(ph_title, title_text,
                   font_name="Arial", size_pt=40, bold=True,
                   color_rgb=KAAR_WHITE)

    # Subtitle / presenter info — find the second text shape by area
    text_shapes = get_text_shapes_by_area(slide, exclude_shape=ph_title)
    presenter_line = "  |  ".join(filter(None, [
        data.get("presenter_name", ""),
        data.get("presenter_id", ""),
        data.get("batch", ""),
    ]))
    if text_shapes and presenter_line.strip("| "):
        force_font(text_shapes[0], presenter_line,
                   font_name="Arial", size_pt=18, bold=False,
                   color_rgb=KAAR_WHITE)

    # Date — third text shape if available
    date_text = data.get("date", "")
    if len(text_shapes) >= 2 and date_text:
        force_font(text_shapes[1], date_text,
                   font_name="Arial", size_pt=14, bold=False,
                   color_rgb=KAAR_WHITE)


# ── AGENDA SLIDE ─────────────────────────────────────────────────────────────

def fill_agenda(slide, data, layout_idx):
    """
    Fill agenda slide. Paired title+body placeholders per item.
    Layout 30 (6 items): title idxs [14,16,18,20,22,24], body [15,17,19,21,23,25]
    Layout 31 (8 items): title idxs [10,12,14,16,18,20,22,24], body [11,13,15,17,19,21,23,25]
    """
    items = data.get("agenda", [])

    # Fill main title (idx=0)
    ph_main = get_ph(slide, 0)
    if ph_main:
        force_font(ph_main, "AGENDA",
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    if layout_idx == 30:   # 6-item layout
        title_idxs = [14, 16, 18, 20, 22, 24]
        body_idxs  = [15, 17, 19, 21, 23, 25]
    else:                  # 8-item layout (layout 31)
        title_idxs = [10, 12, 14, 16, 18, 20, 22, 24]
        body_idxs  = [11, 13, 15, 17, 19, 21, 23, 25]

    for i, (ti, bi) in enumerate(zip(title_idxs, body_idxs)):
        ph_title = get_ph(slide, ti)
        ph_body  = get_ph(slide, bi)
        if i < len(items):
            if ph_title:
                force_font(ph_title, items[i],
                           font_name="Arial", size_pt=13, bold=True,
                           color_rgb=KAAR_DARK)
            if ph_body:
                # Leave body empty for clean look
                force_font(ph_body, "",
                           font_name="Calibri", size_pt=11, bold=False,
                           color_rgb=KAAR_GRAY)
        else:
            # Clear unused slots
            if ph_title:
                force_font(ph_title, "",
                           font_name="Arial", size_pt=13)
            if ph_body:
                force_font(ph_body, "",
                           font_name="Calibri", size_pt=11)


# ── SECTION DIVIDER ─────────────────────────────────────────────────────────

def fill_divider(slide, data):
    """
    Fill section divider slide.
    All splitter layouts use idx=58 "Text Placeholder 35" for main text.
    """
    section_title = data.get("section_title", "SECTION")

    # Main section title — idx=58
    ph = get_ph(slide, 58)
    if ph:
        force_font(ph, section_title.upper(),
                   font_name="Arial", size_pt=36, bold=True,
                   color_rgb=KAAR_WHITE)

    # Optional subtitle — find largest remaining text shape
    subtitle = data.get("section_subtitle", "")
    if subtitle:
        other_shapes = [s for s in slide.shapes
                        if s.has_text_frame
                        and (ph is None or s.shape_id != ph.shape_id)
                        and s.width > 500000]
        if other_shapes:
            other_shapes.sort(key=lambda s: s.width * s.height, reverse=True)
            force_font(other_shapes[0], subtitle,
                       font_name="Arial", size_pt=20, bold=False,
                       color_rgb=KAAR_WHITE)


# ── CONTENT SINGLE (Layout 14) ──────────────────────────────────────────────

def fill_content_single(slide, data):
    """
    Fill single-content slide (Layout 14): idx=0 title + idx=1 content.
    Handles: overview, content_single, dashboard slide types.
    """
    ph_title   = get_ph(slide, 0)
    ph_content = get_ph(slide, 1)

    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    if ph_content:
        # Clear all paragraphs from the content box
        clear_text_body(ph_content.text_frame)

        stype = data.get("slide_type", "overview")

        if stype in ("overview", "content_single"):
            _fill_overview_content(ph_content.text_frame, data)
        elif stype == "dashboard":
            _fill_dashboard_content(ph_content.text_frame, data)
        else:
            # Generic fallback — dump any available text content
            _fill_overview_content(ph_content.text_frame, data)


def _fill_overview_content(tf, data):
    """Fill a text frame with overview-style structured content."""
    # Overview text
    if data.get("overview_text"):
        add_formatted_paragraph(tf, "Overview",
                                font_name="Arial", size_pt=14, bold=True,
                                color_rgb=KAAR_RED)
        add_formatted_paragraph(tf, data["overview_text"],
                                font_name="Calibri", size_pt=12,
                                color_rgb=KAAR_GRAY)

    # SAP Module
    if data.get("sap_module"):
        add_formatted_paragraph(tf, "",
                                font_name="Calibri", size_pt=8)
        add_formatted_paragraph(tf, f"SAP Module: {data['sap_module']}",
                                font_name="Arial", size_pt=13, bold=True,
                                color_rgb=KAAR_RED)
        if data.get("sap_module_description"):
            add_formatted_paragraph(tf, data["sap_module_description"],
                                    font_name="Calibri", size_pt=12,
                                    color_rgb=KAAR_GRAY)

    # Key features
    if data.get("key_features"):
        add_formatted_paragraph(tf, "",
                                font_name="Calibri", size_pt=8)
        add_formatted_paragraph(tf, "Key Features",
                                font_name="Arial", size_pt=13, bold=True,
                                color_rgb=KAAR_RED)
        for f in data["key_features"]:
            add_formatted_paragraph(tf, f"• {f}",
                                    font_name="Calibri", size_pt=12,
                                    color_rgb=KAAR_GRAY)

    # Business benefits
    if data.get("business_benefits"):
        add_formatted_paragraph(tf, "",
                                font_name="Calibri", size_pt=8)
        add_formatted_paragraph(tf, "Business Benefits",
                                font_name="Arial", size_pt=13, bold=True,
                                color_rgb=KAAR_RED)
        for b in data["business_benefits"]:
            add_formatted_paragraph(tf, f"✓ {b}",
                                    font_name="Calibri", size_pt=12,
                                    color_rgb=KAAR_GRAY)


def _fill_dashboard_content(tf, data):
    """Fill a text frame with dashboard-style items."""
    for item in data.get("dashboard_items", []):
        add_formatted_paragraph(tf, f"▶ {item.get('name', '')}",
                                font_name="Arial", size_pt=14, bold=True,
                                color_rgb=KAAR_RED)
        add_formatted_paragraph(tf, item.get("description", ""),
                                font_name="Calibri", size_pt=12,
                                color_rgb=KAAR_GRAY)
        if item.get("tables"):
            add_formatted_paragraph(tf, f"  T-Codes: {item['tables']}",
                                    font_name="Calibri", size_pt=11, italic=True,
                                    color_rgb=KAAR_GRAY)
        add_formatted_paragraph(tf, "",
                                font_name="Calibri", size_pt=6)


# ── CONTENT TWO COLUMNS (Layout 15) ─────────────────────────────────────────

def fill_content_two(slide, data):
    """
    Fill two-column content slide (Layout 15):
    idx=0 title + idx=1 left box + idx=2 right box.
    Handles: components, content_two, comparison.
    """
    ph_title = get_ph(slide, 0)
    ph_left  = get_ph(slide, 1)
    ph_right = get_ph(slide, 2)

    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    stype = data.get("slide_type", "")

    # Handle comparison type specially
    if stype == "comparison":
        _fill_comparison(slide, data, ph_left, ph_right)
        return

    # For components / content_two — split items across two columns
    items = (data.get("components", [])
             or data.get("process_steps", [])
             or data.get("dashboard_items", []))

    mid = max(1, len(items) // 2) if items else 0
    left_items  = items[:mid] if mid > 0 else items
    right_items = items[mid:] if mid > 0 else []

    _fill_column_box(ph_left, left_items)
    _fill_column_box(ph_right, right_items)


def _fill_column_box(ph, item_list):
    """Fill a placeholder box with a list of component/step items."""
    if not ph:
        return
    clear_text_body(ph.text_frame)

    for item in item_list:
        if isinstance(item, dict):
            name = item.get("name") or item.get("step", "")
            desc = item.get("description", "")
            tables = item.get("tables", "")
        else:
            name = str(item)
            desc = ""
            tables = ""

        add_formatted_paragraph(ph.text_frame, name,
                                font_name="Arial", size_pt=14, bold=True,
                                color_rgb=KAAR_RED)
        if desc:
            add_formatted_paragraph(ph.text_frame, desc,
                                    font_name="Calibri", size_pt=12,
                                    color_rgb=KAAR_GRAY)
        if tables:
            add_formatted_paragraph(ph.text_frame, f"T-Codes: {tables}",
                                    font_name="Calibri", size_pt=11, italic=True,
                                    color_rgb=KAAR_GRAY)
        add_formatted_paragraph(ph.text_frame, "",
                                font_name="Calibri", size_pt=6)


def _fill_comparison(slide, data, ph_left, ph_right):
    """Fill a comparison slide with left/right option data."""
    left_data  = data.get("comparison_left", {})
    right_data = data.get("comparison_right", {})

    # For layout 16 (Comparison): idx=1 + idx=3 are headers, idx=2 + idx=4 are content
    # But if we're on layout 15, use the simpler two-box approach
    ph_header_left  = get_ph(slide, 1)
    ph_body_left    = get_ph(slide, 2)
    ph_header_right = get_ph(slide, 3)
    ph_body_right   = get_ph(slide, 4)

    # Layout 16 has 5 placeholders: 0, 1, 2, 3, 4
    if ph_header_right and ph_body_right:
        # Layout 16 — use header+body pairs
        if ph_header_left:
            force_font(ph_header_left, left_data.get("title", "Option A"),
                       font_name="Arial", size_pt=16, bold=True,
                       color_rgb=KAAR_RED)
        if ph_body_left:
            clear_text_body(ph_body_left.text_frame)
            for pt in left_data.get("points", []):
                add_formatted_paragraph(ph_body_left.text_frame, f"• {pt}",
                                        font_name="Calibri", size_pt=12,
                                        color_rgb=KAAR_GRAY)
        if ph_header_right:
            force_font(ph_header_right, right_data.get("title", "Option B"),
                       font_name="Arial", size_pt=16, bold=True,
                       color_rgb=KAAR_RED)
        if ph_body_right:
            clear_text_body(ph_body_right.text_frame)
            for pt in right_data.get("points", []):
                add_formatted_paragraph(ph_body_right.text_frame, f"• {pt}",
                                        font_name="Calibri", size_pt=12,
                                        color_rgb=KAAR_GRAY)
    else:
        # Fallback to simpler two-box layout (layout 15)
        if ph_left:
            clear_text_body(ph_left.text_frame)
            add_formatted_paragraph(ph_left.text_frame,
                                    left_data.get("title", "Option A"),
                                    font_name="Arial", size_pt=16, bold=True,
                                    color_rgb=KAAR_RED)
            for pt in left_data.get("points", []):
                add_formatted_paragraph(ph_left.text_frame, f"• {pt}",
                                        font_name="Calibri", size_pt=12,
                                        color_rgb=KAAR_GRAY)
        if ph_right:
            clear_text_body(ph_right.text_frame)
            add_formatted_paragraph(ph_right.text_frame,
                                    right_data.get("title", "Option B"),
                                    font_name="Arial", size_pt=16, bold=True,
                                    color_rgb=KAAR_RED)
            for pt in right_data.get("points", []):
                add_formatted_paragraph(ph_right.text_frame, f"• {pt}",
                                        font_name="Calibri", size_pt=12,
                                        color_rgb=KAAR_GRAY)


# ── CONTENT THREE COLUMNS (Layout 17) ───────────────────────────────────────

def fill_content_three(slide, data):
    """
    Fill three-column content slide (Layout 17):
    idx=0 title + idx=1 first col + idx=2 second col
    Third column may be at idx=3 or another placeholder.
    """
    ph_title = get_ph(slide, 0)
    ph_col1  = get_ph(slide, 1)
    ph_col2  = get_ph(slide, 2)
    ph_col3  = get_ph(slide, 3)

    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    items = (data.get("components", [])
             or data.get("process_steps", [])
             or data.get("dashboard_items", []))

    # Split into three groups
    third = max(1, len(items) // 3)
    groups = [items[:third], items[third:2*third], items[2*third:]]
    placeholders = [ph_col1, ph_col2, ph_col3]

    for ph, group in zip(placeholders, groups):
        _fill_column_box(ph, group)


# ── PROCESS ARROW (Layout 106) ──────────────────────────────────────────────

def fill_process_arrow(slide, data):
    """
    Fill 4-step arrow process slide (Layout 106):
    idx=0 title, step titles at idx=14,16,18,20, descriptions at idx=15,17,19,21.
    """
    ph_title = get_ph(slide, 0)
    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    steps = data.get("process_steps", [])
    step_title_idxs = [14, 16, 18, 20]
    step_body_idxs  = [15, 17, 19, 21]

    for i, (ti, bi) in enumerate(zip(step_title_idxs, step_body_idxs)):
        if i < len(steps):
            s = steps[i]
            name = s.get("step", "") if isinstance(s, dict) else str(s)
            desc = s.get("description", "") if isinstance(s, dict) else ""
            ph_t = get_ph(slide, ti)
            ph_b = get_ph(slide, bi)
            if ph_t:
                force_font(ph_t, name,
                           font_name="Arial", size_pt=13, bold=True,
                           color_rgb=KAAR_WHITE)
            if ph_b:
                force_font(ph_b, desc,
                           font_name="Calibri", size_pt=11, bold=False,
                           color_rgb=KAAR_WHITE)
        else:
            # Clear unused arrow steps
            ph_t = get_ph(slide, ti)
            ph_b = get_ph(slide, bi)
            if ph_t:
                force_font(ph_t, "", font_name="Arial", size_pt=13)
            if ph_b:
                force_font(ph_b, "", font_name="Calibri", size_pt=11)


# ── PROCESS FLOW (Layout 135) ───────────────────────────────────────────────

def fill_process_flow(slide, data):
    """
    Fill 5+ step flow diagram (Layout 135):
    idx=0 title + step content at idx=51,53,64,65,66,67,68.
    """
    ph_title = get_ph(slide, 0)
    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    steps = data.get("process_steps", [])
    step_idxs = [51, 53, 64, 65, 66, 67, 68]

    for i, si in enumerate(step_idxs):
        ph = get_ph(slide, si)
        if ph:
            if i < len(steps):
                s = steps[i]
                if isinstance(s, dict):
                    text = s.get("step", "")
                    desc = s.get("description", "")
                    text = f"{text} — {desc}" if desc else text
                else:
                    text = str(s)
                force_font(ph, text[:120],
                           font_name="Calibri", size_pt=12, bold=False,
                           color_rgb=KAAR_DARK)
            else:
                force_font(ph, "", font_name="Calibri", size_pt=12)


# ── PROCESS CYCLE (Layout 84) ───────────────────────────────────────────────

def fill_process_cycle(slide, data):
    """
    Fill circular cycle slide (Layout 84):
    idx=0 title + step titles at idx=51,56,58,59,60 + descriptions at idx=55,57,61,62,63.
    """
    ph_title = get_ph(slide, 0)
    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    steps = data.get("process_steps", []) or data.get("cycle_steps", [])
    step_title_idxs = [51, 56, 58, 59, 60]
    step_desc_idxs  = [55, 57, 61, 62, 63]

    for i, (ti, di) in enumerate(zip(step_title_idxs, step_desc_idxs)):
        ph_t = get_ph(slide, ti)
        ph_d = get_ph(slide, di)
        if i < len(steps):
            s = steps[i]
            if isinstance(s, dict):
                name = s.get("step", "")
                desc = s.get("description", "")
            else:
                name = str(s)
                desc = ""
            if ph_t:
                force_font(ph_t, name,
                           font_name="Arial", size_pt=13, bold=True,
                           color_rgb=KAAR_DARK)
            if ph_d:
                force_font(ph_d, desc,
                           font_name="Calibri", size_pt=11, bold=False,
                           color_rgb=KAAR_GRAY)
        else:
            if ph_t:
                force_font(ph_t, "", font_name="Arial", size_pt=13)
            if ph_d:
                force_font(ph_d, "", font_name="Calibri", size_pt=11)


# ── TIMELINE (Layout 59) ────────────────────────────────────────────────────

def fill_timeline(slide, data):
    """
    Fill horizontal timeline slide (Layout 59):
    idx=0 title + idx=10..18 (9 text placeholders for step labels).
    """
    ph_title = get_ph(slide, 0)
    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    steps = data.get("timeline_steps", []) or data.get("process_steps", [])
    timeline_idxs = list(range(10, 19))  # idx 10..18

    for i, ti in enumerate(timeline_idxs):
        ph = get_ph(slide, ti)
        if ph:
            if i < len(steps):
                s = steps[i]
                text = s.get("step", s) if isinstance(s, dict) else str(s)
                force_font(ph, text,
                           font_name="Arial", size_pt=12, bold=True,
                           color_rgb=KAAR_DARK)
            else:
                force_font(ph, "", font_name="Arial", size_pt=12)


# ── POINTER 4 (Layout 96) ───────────────────────────────────────────────────

def fill_pointer_4(slide, data):
    """
    Fill 4-pointer slide (Layout 96):
    idx=0 title + idx=14..17 descriptions + idx=18..21 number labels.
    Icon placeholders idx=10..13 are skipped.
    """
    ph_title = get_ph(slide, 0)
    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    items = (data.get("key_features", [])
             or data.get("components", [])
             or data.get("process_steps", []))
    desc_idxs   = [14, 15, 16, 17]
    number_idxs = [18, 19, 20, 21]

    for i in range(4):
        ph_desc = get_ph(slide, desc_idxs[i])
        ph_num  = get_ph(slide, number_idxs[i])

        if i < len(items):
            item = items[i]
            if isinstance(item, dict):
                text = item.get("name", "") or item.get("step", "")
                desc = item.get("description", "")
                display = f"{text}\n{desc}" if desc else text
            else:
                display = str(item)

            if ph_desc:
                force_font(ph_desc, display,
                           font_name="Calibri", size_pt=12, bold=False,
                           color_rgb=KAAR_GRAY)
            if ph_num:
                force_font(ph_num, str(i + 1),
                           font_name="Arial", size_pt=20, bold=True,
                           color_rgb=KAAR_RED)
        else:
            if ph_desc:
                force_font(ph_desc, "", font_name="Calibri", size_pt=12)
            if ph_num:
                force_font(ph_num, "", font_name="Arial", size_pt=20)


# ── CIRCLE 4 (Layout 74) ────────────────────────────────────────────────────

def fill_circle_4(slide, data):
    """
    Fill 4-item circle diagram (Layout 74):
    idx=0 title + idx=10..13 labels + idx=14..17 content.
    """
    ph_title = get_ph(slide, 0)
    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    items = (data.get("components", [])
             or data.get("key_features", [])
             or data.get("process_steps", []))
    label_idxs   = [10, 11, 12, 13]
    content_idxs = [14, 15, 16, 17]

    for i in range(4):
        ph_label   = get_ph(slide, label_idxs[i])
        ph_content = get_ph(slide, content_idxs[i])

        if i < len(items):
            item = items[i]
            if isinstance(item, dict):
                name = item.get("name", "") or item.get("step", "")
                desc = item.get("description", "")
            else:
                name = str(item)
                desc = ""

            if ph_label:
                force_font(ph_label, name,
                           font_name="Arial", size_pt=13, bold=True,
                           color_rgb=KAAR_DARK)
            if ph_content:
                force_font(ph_content, desc,
                           font_name="Calibri", size_pt=11, bold=False,
                           color_rgb=KAAR_GRAY)
        else:
            if ph_label:
                force_font(ph_label, "", font_name="Arial", size_pt=13)
            if ph_content:
                force_font(ph_content, "", font_name="Calibri", size_pt=11)


# ── QUOTE (Layout 45) ───────────────────────────────────────────────────────

def fill_quote(slide, data):
    """Fill quote slide with main quote text and attribution."""
    # Try to find placeholders — quote slides vary
    text_shapes = get_text_shapes_by_area(slide)

    quote_text = data.get("quote_text", "")
    attribution = data.get("quote_attribution", "")

    if text_shapes:
        force_font(text_shapes[0], quote_text,
                   font_name="Arial", size_pt=24, bold=True,
                   color_rgb=KAAR_DARK)
    if len(text_shapes) >= 2 and attribution:
        force_font(text_shapes[1], attribution,
                   font_name="Calibri", size_pt=14, italic=True,
                   color_rgb=KAAR_GRAY)


# ── TABLE (Layout 116) ──────────────────────────────────────────────────────

def fill_table(slide, data):
    """
    Fill table slide. Layout 116 has many cell placeholders.
    Fallback to content_single if structure doesn't match.
    """
    ph_title = get_ph(slide, 0)
    if ph_title:
        force_font(ph_title, data.get("title", ""),
                   font_name="Arial", size_pt=28, bold=True,
                   color_rgb=KAAR_RED)

    headers = data.get("table_headers", [])
    rows    = data.get("table_rows", [])

    # Try to fill cell placeholders idx=10..50
    cell_idx = 10
    # Fill headers first
    for h in headers:
        ph = get_ph(slide, cell_idx)
        if ph:
            force_font(ph, h,
                       font_name="Arial", size_pt=12, bold=True,
                       color_rgb=KAAR_WHITE)
        cell_idx += 1

    # Fill rows
    for row in rows:
        for val in row:
            ph = get_ph(slide, cell_idx)
            if ph:
                force_font(ph, str(val),
                           font_name="Calibri", size_pt=11, bold=False,
                           color_rgb=KAAR_DARK)
            cell_idx += 1


# ── THANK YOU (Layout 174) ──────────────────────────────────────────────────

def fill_thankyou(slide, data):
    """
    Fill thank you slide (Layout 174):
    idx=0 for contact details. TextBox "Thank You" already present (non-placeholder).
    """
    thank_you_data = data.get("thank_you", {})
    contact = thank_you_data.get("contact", "")
    msg     = thank_you_data.get("message", "Thank You")

    ph_contact = get_ph(slide, 0)
    if ph_contact and contact:
        force_font(ph_contact, contact,
                   font_name="Calibri", size_pt=14, bold=False,
                   color_rgb=KAAR_DARK)

    # The "Thank You" text is in a NON-PLACEHOLDER TextBox (e.g. "TextBox 7")
    # It already says "Thank You" — only overwrite if custom message provided
    if msg and msg != "Thank You":
        for shape in slide.shapes:
            if shape.has_text_frame and "Thank You" in shape.text_frame.text:
                force_font(shape, msg,
                           font_name="Arial", size_pt=44, bold=True,
                           color_rgb=KAAR_DARK)
                break


# ═════════════════════════════════════════════════════════════════════════════
#  GEMINI CLIENT
# ═════════════════════════════════════════════════════════════════════════════

class GeminiClient:
    """Wrapper around Google Generative AI SDK (Gemini 2.0 Flash)."""

    MODEL = "gemini-2.0-flash"

    def __init__(self, api_key: Optional[str] = None):
        key = api_key or os.getenv("GEMINI_API_KEY")
        if not key:
            raise ValueError(
                "GEMINI_API_KEY not found. Set it in .env or pass --api-key."
            )
        genai.configure(api_key=key)
        self.model = genai.GenerativeModel(
            model_name=self.MODEL,
            generation_config=genai.GenerationConfig(
                response_mime_type="application/json",
                temperature=0.3,
                max_output_tokens=8192,
            )
        )

    def generate_json(self, prompt: str) -> dict:
        """Generate JSON content from Gemini with retry logic."""
        for attempt in range(3):
            try:
                response = self.model.generate_content(prompt)
                raw = response.text.strip()
                # Strip code fences if present
                if "```" in raw:
                    for part in raw.split("```"):
                        cleaned = part.strip().lstrip("json").strip()
                        if cleaned.startswith("{"):
                            raw = cleaned
                            break
                # Extract JSON object
                start = raw.find("{")
                end   = raw.rfind("}") + 1
                if start != -1 and end > start:
                    raw = raw[start:end]
                result = json.loads(raw)
                print(f"[Gemini] Parsed OK on attempt {attempt + 1}")
                return result
            except json.JSONDecodeError as e:
                print(f"[Gemini] Parse failed attempt {attempt + 1}: {e}")
                if attempt < 2:
                    prompt = ("IMPORTANT: Your last response was not valid JSON. "
                              + prompt)
        raise ValueError("Gemini failed to return valid JSON after 3 attempts.")

    def generate_text(self, prompt: str) -> str:
        """Return plain text from Gemini."""
        response = self.model.generate_content(prompt)
        return response.text.strip()


# ═════════════════════════════════════════════════════════════════════════════
#  GEMINI PROMPT BUILDER
# ═════════════════════════════════════════════════════════════════════════════

def build_gemini_prompt(user_topic: str, context: str = "") -> str:
    """Build the structured prompt for Gemini to produce presentation JSON."""
    return f"""
SYSTEM: You are a senior SAP consultant at Kaar Technologies.
OUTPUT RULE: Respond with ONLY valid JSON. No markdown. No code fences.
Start with {{ and end with }}. No text before or after.

USER REQUEST: {user_topic}
CONTEXT: {context or "None"}

Generate a complete presentation JSON using EXACTLY this structure.
Replace all <placeholder> values with real, accurate content.

SLIDE TYPE GUIDE — pick the best type for each slide:
  overview       → intro slide with overview_text + key_features + sap_module
  content_single → single column text/bullets
  content_two    → two parallel columns (equal importance)
  components     → list of named components with descriptions + t-codes
  process        → sequential steps (use process_steps array)
  dashboard      → monitoring metrics with dashboard_items array
  comparison     → side-by-side two options
  timeline       → time-ordered steps (use timeline_steps array, max 9)
  cycle          → repeating circular process
  pointer_4      → exactly 4 key highlighted items
  circle_4       → exactly 4 connected concepts in circle
  quote          → one key stat or statement (use quote_text)
  table          → tabular data (use table_headers + table_rows)

REQUIRED CONTENT DEPTH:
  - overview slides: minimum 6 key_features, 5 business_benefits
  - process slides: minimum 5 process_steps, each with a full sentence description
  - components slides: minimum 4 components, each with real SAP t-codes
  - dashboard slides: minimum 5 dashboard_items with real SAP t-codes
  - agenda: exactly 7 items (one per major topic)
  - Use REAL SAP transaction codes (SE10, STMS, SM50, SM37, DB02, etc.)
  - Use REAL SAP table names where relevant (VBAK, PA0001, EKKO, etc.)

{{
  "presentation_title": "<title matching user request>",
  "subtitle": "<one-line subtitle>",
  "presenter_name": "THISO VALLABADASS K",
  "presenter_id": "AID-902015",
  "batch": "GE-26 Batch-1",
  "date": "07-06-2026",
  "agenda": [
    "<topic 1>", "<topic 2>", "<topic 3>",
    "<topic 4>", "<topic 5>", "<topic 6>", "<topic 7>"
  ],
  "sections": [
    {{
      "section_title": "<SECTION NAME IN CAPITALS>",
      "section_subtitle": "<subtitle or context>",
      "slides": [
        {{
          "slide_type": "overview",
          "title": "<slide title>",
          "overview_text": "<2-3 sentence description>",
          "sap_module": "<e.g. SAP NetWeaver>",
          "sap_module_description": "<one sentence per module>",
          "key_features": ["<f1>","<f2>","<f3>","<f4>","<f5>","<f6>"],
          "business_benefits": ["<b1>","<b2>","<b3>","<b4>","<b5>"]
        }},
        {{
          "slide_type": "components",
          "title": "<slide title>",
          "components": [
            {{"name":"<name>","description":"<desc>","tables":"<SAP t-codes>"}},
            {{"name":"<name>","description":"<desc>","tables":"<SAP t-codes>"}},
            {{"name":"<name>","description":"<desc>","tables":"<SAP t-codes>"}},
            {{"name":"<name>","description":"<desc>","tables":"<SAP t-codes>"}}
          ]
        }},
        {{
          "slide_type": "process",
          "title": "<PROCESS NAME>",
          "process_steps": [
            {{"step":"<name>","description":"<one sentence>"}},
            {{"step":"<name>","description":"<one sentence>"}},
            {{"step":"<name>","description":"<one sentence>"}},
            {{"step":"<name>","description":"<one sentence>"}},
            {{"step":"<name>","description":"<one sentence>"}},
            {{"step":"<name>","description":"<one sentence>"}}
          ]
        }},
        {{
          "slide_type": "dashboard",
          "title": "<DASHBOARD TITLE>",
          "dashboard_items": [
            {{"name":"<metric>","description":"<what it shows>","tables":"<t-codes>"}},
            {{"name":"<metric>","description":"<what it shows>","tables":"<t-codes>"}},
            {{"name":"<metric>","description":"<what it shows>","tables":"<t-codes>"}},
            {{"name":"<metric>","description":"<what it shows>","tables":"<t-codes>"}},
            {{"name":"<metric>","description":"<what it shows>","tables":"<t-codes>"}}
          ]
        }}
      ]
    }}
  ],
  "thank_you": {{
    "message": "Thank You",
    "contact": "thiso.k@kaartech.com"
  }}
}}

RULES:
- Output ONLY the JSON. Nothing before {{. Nothing after }}.
- Use real SAP transaction codes (SE10, STMS, SM50, SM37, DB02, etc.)
- Use professional enterprise language
- Every array must have at least the minimum items shown above
- Each section should have 3-5 slides with varied slide_types

Generate ONLY the JSON. Start with {{.
"""


# ═════════════════════════════════════════════════════════════════════════════
#  BUILD FUNCTION — orchestrates the entire presentation
# ═════════════════════════════════════════════════════════════════════════════

def build(prs_path, content, output_path, debug=False):
    """
    Build the full presentation from a content dict.

    Args:
        prs_path:    Path to reference template .pptx
        content:     dict matching the JSON schema
        output_path: where to save the output .pptx
        debug:       if True, print placeholder info for each slide

    Returns:
        str: output file path
    """
    prs = Presentation(str(prs_path))

    # ── Remove all existing slides (keep masters/layouts/theme) ──────────────
    slide_id_list = prs.slides._sldIdLst
    for sld_id in list(slide_id_list):
        try:
            prs.part.drop_rel(sld_id.rId)
        except Exception:
            pass
        slide_id_list.remove(sld_id)

    title    = content.get("presentation_title", "KaarTech Presentation")
    agenda   = content.get("agenda", [])
    sections = content.get("sections", [])

    # ── 1. Cover slide ───────────────────────────────────────────────────────
    cover_layout_key = content.get("cover_style", "cover")
    cover_layout_idx = LAYOUT_MAP.get(cover_layout_key, LAYOUT_MAP["cover"])
    cover_slide = add_slide_from_layout(prs, cover_layout_idx)
    fill_cover(cover_slide, {
        "title":          title,
        "presenter_name": content.get("presenter_name", ""),
        "presenter_id":   content.get("presenter_id", ""),
        "batch":          content.get("batch", ""),
        "date":           content.get("date", ""),
    })
    print(f"[Build] Cover: {title}")

    if debug:
        _debug_slide(cover_slide, "Cover")

    # ── 2. Agenda slide ──────────────────────────────────────────────────────
    agenda_layout_idx = 30 if len(agenda) <= 6 else 31
    agenda_slide = add_slide_from_layout(prs, agenda_layout_idx)
    fill_agenda(agenda_slide, content, agenda_layout_idx)
    print(f"[Build] Agenda: {len(agenda)} items (layout {agenda_layout_idx})")

    if debug:
        _debug_slide(agenda_slide, "Agenda")

    # ── 3. Sections ──────────────────────────────────────────────────────────
    for sec_idx, sec in enumerate(sections):
        sec_title = sec.get("section_title", f"Section {sec_idx + 1}")

        # Section divider
        div_layout = LAYOUT_MAP.get(
            sec.get("divider_style", "divider"),
            LAYOUT_MAP["divider"]
        )
        div_slide = add_slide_from_layout(prs, div_layout)
        fill_divider(div_slide, sec)
        print(f"[Build] Section divider: {sec_title}")

        if debug:
            _debug_slide(div_slide, f"Divider: {sec_title}")

        # Section slides
        for slide_data in sec.get("slides", []):
            layout_idx = choose_layout(slide_data, agenda)
            slide = add_slide_from_layout(prs, layout_idx)
            stype = slide_data.get("slide_type", "overview")
            slide_title = slide_data.get("title", "")
            print(f"  [Slide] {stype} (layout {layout_idx}): {slide_title}")

            # ── Route to the correct fill function ───────────────────────────
            if stype in ("overview", "content_single"):
                fill_content_single(slide, slide_data)

            elif stype == "dashboard":
                fill_content_single(slide, slide_data)

            elif stype in ("components", "content_two"):
                if layout_idx == 17:
                    fill_content_three(slide, slide_data)
                else:
                    fill_content_two(slide, slide_data)

            elif stype == "comparison":
                fill_content_two(slide, slide_data)

            elif stype == "process":
                steps = slide_data.get("process_steps", [])
                if slide_data.get("is_cycle"):
                    fill_process_cycle(slide, slide_data)
                elif len(steps) <= 4:
                    fill_process_arrow(slide, slide_data)
                else:
                    fill_process_flow(slide, slide_data)

            elif stype == "cycle":
                fill_process_cycle(slide, slide_data)

            elif stype == "timeline":
                fill_timeline(slide, slide_data)

            elif stype == "pointer_4":
                fill_pointer_4(slide, slide_data)

            elif stype == "circle_4":
                fill_circle_4(slide, slide_data)

            elif stype == "quote":
                fill_quote(slide, slide_data)

            elif stype == "table":
                fill_table(slide, slide_data)

            elif stype == "content_three":
                fill_content_three(slide, slide_data)

            else:
                # Default fallback to content_single
                fill_content_single(slide, slide_data)

            if debug:
                _debug_slide(slide, f"{stype}: {slide_title}")

    # ── 4. Thank You slide ───────────────────────────────────────────────────
    ty_layout = LAYOUT_MAP.get(
        content.get("thankyou_style", "thankyou"),
        LAYOUT_MAP["thankyou"]
    )
    ty_slide = add_slide_from_layout(prs, ty_layout)
    fill_thankyou(ty_slide, content)
    print("[Build] Thank you slide")

    if debug:
        _debug_slide(ty_slide, "Thank You")

    # ── Save ─────────────────────────────────────────────────────────────────
    prs.save(str(output_path))
    print(f"\n[Done] Saved: {output_path}")
    return str(output_path)


def _debug_slide(slide, label):
    """Print all shape names, sizes, and text previews for debugging."""
    clean_label = label.encode('ascii', errors='replace').decode('ascii')
    print(f"\n  [DEBUG] Slide shapes — {clean_label}")
    for s in slide.shapes:
        is_ph = hasattr(s, 'placeholder_format') and s.placeholder_format is not None
        ph_idx = s.placeholder_format.idx if is_ph else "-"
        if s.has_text_frame:
            txt = s.text_frame.text[:40].replace('\n', ' ')
            txt = txt.encode('ascii', errors='replace').decode('ascii')
        else:
            txt = "(no text)"
        name = s.name.encode('ascii', errors='replace').decode('ascii')
        print(f"    ph={str(ph_idx):>3s}  name={name!r:30s}  "
              f"w={s.width/914400:.1f}\"  h={s.height/914400:.1f}\"  "
              f"txt={txt!r}")


# ═════════════════════════════════════════════════════════════════════════════
#  INSPECT & EDIT EXISTING PPTX
# ═════════════════════════════════════════════════════════════════════════════

def inspect_pptx(pptx_path: str) -> dict:
    """
    Open any .pptx file and return its full structure as a JSON-serializable dict.
    This is how Claude learns what's in an existing presentation before editing it.
    """
    prs = Presentation(pptx_path)
    result = {
        "file": os.path.basename(pptx_path),
        "total_slides": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "layout_name": slide.slide_layout.name,
            "shapes": []
        }

        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "placeholder_idx": None,
                "has_text": shape.has_text_frame,
                "text_preview": "",
                "full_text": "",
                "font_info": {}
            }

            # Get placeholder index if applicable
            if shape.is_placeholder:
                try:
                    shape_info["placeholder_idx"] = shape.placeholder_format.idx
                except Exception:
                    pass

            # Get text content and font info
            if shape.has_text_frame:
                full_text = shape.text_frame.text
                shape_info["full_text"] = full_text
                shape_info["text_preview"] = full_text[:80].replace('\n', ' ')

                # Get font from first run of first paragraph
                try:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        font = run.font
                        shape_info["font_info"] = {
                            "size_pt": round(font.size / 12700) if font.size else None,
                            "bold": font.bold,
                            "font_name": font.name,
                            "color_hex": (
                                f"#{font.color.rgb}" if font.color and font.color.type else None
                            )
                        }
                except Exception:
                    pass

            slide_info["shapes"].append(shape_info)

        result["slides"].append(slide_info)

    return result


def edit_pptx(pptx_path: str, edits: list, output_path: str,
              gemini_client=None) -> str:
    """
    Open an existing .pptx, apply a list of edits, save a COPY to output_path.
    The original file is NEVER modified.

    edits is a list of edit operation dicts. Each operation has:
      {
        "operation": "update_text" | "update_slide_title" | "add_bullet" |
                     "replace_text" | "add_slide_after" | "delete_slide" |
                     "update_section_title" | "generate_slide_content" |
                     "replace_text_all_slides",
        "slide_number": <int, 1-based>,
        ... operation-specific fields ...
      }

    Returns: output_path (str) — path to the saved edited copy
    """
    # ── Safety: never modify original ─────────────────────────────────────────
    shutil.copy2(pptx_path, output_path)
    prs = Presentation(output_path)

    def get_slide(slide_number):
        """Return 0-based slide by 1-based number."""
        idx = slide_number - 1
        if idx < 0 or idx >= len(prs.slides):
            raise ValueError(f"Slide {slide_number} does not exist (total: {len(prs.slides)})")
        return prs.slides[idx]

    def find_shape(slide, shape_name=None, placeholder_idx=None):
        """Find a shape by name (partial match) or placeholder idx."""
        if placeholder_idx is not None:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == placeholder_idx:
                    return shape
        if shape_name:
            name_lower = shape_name.lower()
            for shape in slide.shapes:
                if name_lower in shape.name.lower():
                    return shape
        return None

    def set_shape_text(shape, new_text, font_name=None, font_size=None,
                       bold=None, color_hex=None):
        """
        Replace all text in a shape. Set font explicitly at XML level
        to override any inherited theme fonts.
        """
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        txBody = tf._txBody
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        # Remove all existing <a:p> elements
        for p in list(txBody.findall(f'{{{ns}}}p')):
            txBody.remove(p)

        # Create a single paragraph with explicit run formatting
        p_elem = etree.SubElement(txBody, f'{{{ns}}}p')
        r_elem = etree.SubElement(p_elem, f'{{{ns}}}r')
        rPr_attrib = {'lang': 'en-US', 'dirty': '0'}
        rPr = etree.SubElement(r_elem, f'{{{ns}}}rPr', attrib=rPr_attrib)

        if bold is not None:
            rPr.set('b', '1' if bold else '0')
        if font_size is not None:
            rPr.set('sz', str(int(font_size * 100)))
        if font_name:
            etree.SubElement(rPr, f'{{{ns}}}latin', attrib={'typeface': font_name})
        if color_hex:
            hex_clean = color_hex.lstrip('#')
            solidFill = etree.SubElement(rPr, f'{{{ns}}}solidFill')
            etree.SubElement(solidFill, f'{{{ns}}}srgbClr', attrib={'val': hex_clean})

        t_elem = etree.SubElement(r_elem, f'{{{ns}}}t')
        t_elem.text = new_text

    def append_bullet(shape, bullet_text, font_size=13, bold=False,
                      color_hex="585858"):
        """Append a paragraph/bullet to a shape's text frame."""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        txBody = tf._txBody
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        p_elem = etree.SubElement(txBody, f'{{{ns}}}p')
        r_elem = etree.SubElement(p_elem, f'{{{ns}}}r')
        rPr = etree.SubElement(r_elem, f'{{{ns}}}rPr',
                                attrib={'lang': 'en-US', 'dirty': '0'})
        rPr.set('b', '1' if bold else '0')
        rPr.set('sz', str(int(font_size * 100)))
        hex_clean = color_hex.lstrip('#')
        solidFill = etree.SubElement(rPr, f'{{{ns}}}solidFill')
        etree.SubElement(solidFill, f'{{{ns}}}srgbClr', attrib={'val': hex_clean})
        t_elem = etree.SubElement(r_elem, f'{{{ns}}}t')
        t_elem.text = bullet_text

    def replace_in_shape(shape, find_str, replace_str):
        """Replace text in all runs of a shape."""
        if not shape.has_text_frame:
            return
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if find_str in run.text:
                    run.text = run.text.replace(find_str, replace_str)

    def _add_slide_after(prs, after_slide_number, layout_name, content):
        """
        Add a new slide after the given position using layout by name.
        Inserts at position after_slide_number (1-based).
        """
        # Find the layout by name
        target_layout = None
        for layout in prs.slide_layouts:
            if layout_name.lower() in layout.name.lower():
                target_layout = layout
                break
        if not target_layout:
            # Fallback: use first content layout
            for layout in prs.slide_layouts:
                if "content" in layout.name.lower():
                    target_layout = layout
                    break
        if not target_layout:
            target_layout = prs.slide_layouts[0]

        new_slide = prs.slides.add_slide(target_layout)

        # Move slide to correct position
        xml_slides = prs.slides._sldIdLst
        last = xml_slides[-1]  # the newly added slide
        xml_slides.remove(last)
        insert_pos = min(after_slide_number, len(xml_slides))
        xml_slides.insert(insert_pos, last)

        # Fill content
        title_text = content.get("title", "")
        body_text  = content.get("body", "")
        bullets    = content.get("bullets", [])

        # Try to set title (placeholder idx=0)
        for shape in new_slide.placeholders:
            if shape.placeholder_format.idx == 0:
                set_shape_text(shape, title_text,
                               font_name="Arial", font_size=28,
                               bold=True, color_hex="C0202B")
                break

        # Try to set body (placeholder idx=1)
        for shape in new_slide.placeholders:
            if shape.placeholder_format.idx == 1:
                lines = []
                if body_text:
                    lines.append(body_text)
                lines.extend(bullets)
                set_shape_text(shape, "\n".join(lines),
                               font_name="Calibri", font_size=13,
                               bold=False, color_hex="585858")
                break

        return new_slide

    def _delete_slide(prs, slide_number):
        """Remove a slide by 1-based number."""
        xml_slides = prs.slides._sldIdLst
        idx = slide_number - 1
        if 0 <= idx < len(xml_slides):
            xml_slides.remove(xml_slides[idx])

    # ── Separate delete operations and sort in reverse to avoid index drift ───
    delete_ops = [e for e in edits if e.get("operation") == "delete_slide"]
    other_ops  = [e for e in edits if e.get("operation") != "delete_slide"]

    # ── Apply non-delete operations first ─────────────────────────────────────
    for edit in other_ops:
        op = edit.get("operation", "")
        snum = edit.get("slide_number")

        try:
            if op == "update_text":
                slide = get_slide(snum)
                shape = find_shape(slide,
                                   shape_name=edit.get("shape_name"),
                                   placeholder_idx=edit.get("placeholder_idx"))
                if shape:
                    set_shape_text(shape, edit["new_text"],
                                   font_name=edit.get("font_name"),
                                   font_size=edit.get("font_size"),
                                   bold=edit.get("bold"),
                                   color_hex=edit.get("color_hex"))
                    print(f"  [Edit] update_text slide={snum} shape={edit.get('shape_name')}")
                else:
                    print(f"  [WARN] Shape not found on slide {snum}: {edit.get('shape_name')}")

            elif op == "update_slide_title":
                slide = get_slide(snum)
                shape = find_shape(slide, placeholder_idx=0)
                if not shape:
                    shape = find_shape(slide, shape_name="Title")
                if shape:
                    set_shape_text(shape, edit["new_title"],
                                   font_name="Arial", font_size=28,
                                   bold=True, color_hex="C0202B")
                    print(f"  [Edit] update_slide_title slide={snum}: {edit['new_title']}")

            elif op == "update_section_title":
                slide = get_slide(snum)
                text_shapes = [s for s in slide.shapes if s.has_text_frame]
                if text_shapes:
                    text_shapes.sort(key=lambda s: s.width * s.height, reverse=True)
                    set_shape_text(text_shapes[0], edit["new_title"].upper(),
                                   font_name="Arial", font_size=36,
                                   bold=True, color_hex="FFFFFF")
                    print(f"  [Edit] update_section_title slide={snum}: {edit['new_title']}")

            elif op == "replace_text":
                slide = get_slide(snum)
                for shape in slide.shapes:
                    replace_in_shape(shape, edit["find"], edit["replace"])
                print(f"  [Edit] replace_text slide={snum}: '{edit['find']}' -> '{edit['replace']}'")

            elif op == "replace_text_all_slides":
                for slide in prs.slides:
                    for shape in slide.shapes:
                        replace_in_shape(shape, edit["find"], edit["replace"])
                print(f"  [Edit] replace_text_all_slides: '{edit['find']}' -> '{edit['replace']}'")

            elif op == "add_bullet":
                slide = get_slide(snum)
                shape = find_shape(slide, shape_name=edit.get("shape_name"))
                if shape:
                    append_bullet(shape, edit["bullet_text"],
                                  font_size=edit.get("font_size", 13),
                                  bold=edit.get("bold", False),
                                  color_hex=edit.get("color_hex", "585858"))
                    print(f"  [Edit] add_bullet slide={snum}")

            elif op == "add_slide_after":
                _add_slide_after(prs, snum,
                                 edit.get("layout_name", "Content Slide_3 - Title and Content"),
                                 edit.get("content", {}))
                print(f"  [Edit] add_slide_after slide={snum}")

            elif op == "generate_slide_content":
                if gemini_client is None:
                    print(f"  [WARN] generate_slide_content requires Gemini client")
                    continue
                slide = get_slide(snum)
                gen_prompt = (
                    f"Generate professional slide content for a KaarTech SAP presentation. "
                    f"Request: {edit.get('prompt', '')}. "
                    f"Return ONLY the text content, no markdown, no headers. "
                    f"Maximum 150 words."
                )
                generated = gemini_client.generate_text(gen_prompt)
                shape = find_shape(slide,
                                   shape_name=edit.get("target_shape_name"),
                                   placeholder_idx=edit.get("placeholder_idx"))
                if shape:
                    set_shape_text(shape, generated,
                                   font_name="Calibri", font_size=13,
                                   bold=False, color_hex="585858")
                    print(f"  [Edit] generate_slide_content slide={snum} (Gemini)")

        except Exception as e:
            print(f"  [ERROR] op={op} slide={snum}: {e}")

    # ── Apply deletes in reverse order ────────────────────────────────────────
    for edit in sorted(delete_ops, key=lambda e: e.get("slide_number", 0), reverse=True):
        snum = edit.get("slide_number")
        try:
            _delete_slide(prs, snum)
            print(f"  [Edit] delete_slide slide={snum}")
        except Exception as e:
            print(f"  [ERROR] delete_slide slide={snum}: {e}")

    prs.save(output_path)
    print(f"[Edit] Saved edited copy to: {output_path}")
    return output_path


# ═════════════════════════════════════════════════════════════════════════════
#  CLI — ARGUMENT PARSING & MAIN
# ═════════════════════════════════════════════════════════════════════════════

def parse_args():
    """Parse command-line arguments."""
    p = argparse.ArgumentParser(
        description="KaarTech Corporate PPT Generator v2.0 — powered by Gemini AI"
    )
    p.add_argument("--topic", "-t",
                   help="Topic / high-level description of the presentation")
    p.add_argument("--prompt-file", "-p",
                   help="Path to a .txt file containing the detailed prompt")
    p.add_argument("--json-file", "-j",
                   help="Path to a pre-generated content JSON file (skip Gemini)")
    p.add_argument("--output", "-o",
                   help="Output filename (without path); saved to output/")
    p.add_argument("--api-key", "-k",
                   help="Gemini API key (overrides .env)")
    p.add_argument("--template", default=str(TEMPLATE_PATH),
                   help=f"Path to KaarTech template .pptx (default: {TEMPLATE_PATH})")
    p.add_argument("--context", "-c", default="",
                   help="Additional context injected into the Gemini prompt")
    p.add_argument("--debug", action="store_true",
                   help="Print placeholder names and sizes for each slide")
    p.add_argument("--edit-file", "-e",
                   help="Path to existing .pptx to edit (instead of generating new)")
    p.add_argument("--edits-json", "-x",
                   help="JSON string or path to .json file containing edit operations list")
    p.add_argument("--inspect", action="store_true",
                   help="Inspect a .pptx file (use with --edit-file) and print slide structure")
    return p.parse_args()


def main():
    """CLI entry point — ties together Gemini content generation and PPT build."""
    args = parse_args()

    # ── INSPECT MODE ─────────────────────────────────────────────────────────
    if args.inspect and args.edit_file:
        import json as _json
        structure = inspect_pptx(args.edit_file)
        output_json = OUTPUT_DIR / "inspected_structure.json"
        output_json.write_text(_json.dumps(structure, indent=2), encoding="utf-8")
        print(f"[Inspect] {structure['total_slides']} slides found in {structure['file']}")
        print(f"[Inspect] Structure saved to: {output_json}")
        for s in structure["slides"]:
            shapes_summary = ", ".join(
                f"{sh['name']}={sh['text_preview'][:30]!r}"
                for sh in s["shapes"] if sh["has_text"]
            )
            print(f"  Slide {s['slide_number']:2d} [{s['layout_name']}] {shapes_summary}")
        return

    # ── EDIT MODE ─────────────────────────────────────────────────────────────
    if args.edit_file:
        import json as _json
        edit_file_path = Path(args.edit_file)
        if not edit_file_path.exists():
            print(f"[Error] File not found: {edit_file_path}")
            sys.exit(1)

        # Parse edits from --edits-json (JSON string or path)
        if not args.edits_json:
            print("[Error] --edits-json is required with --edit-file")
            sys.exit(1)

        edits_input = args.edits_json.strip()
        if Path(edits_input).exists():
            with open(edits_input, "r", encoding="utf-8-sig") as f:
                edits = _json.load(f)
        else:
            edits = _json.loads(edits_input)

        # Set output filename
        out_name = args.output or f"edited_{edit_file_path.stem}.pptx"
        output_path = OUTPUT_DIR / out_name

        # Optional Gemini for generate_slide_content operations
        needs_gemini = any(e.get("operation") == "generate_slide_content" for e in edits)
        gemini = None
        if needs_gemini:
            gemini = GeminiClient(api_key=args.api_key)

        print(f"[Edit] Opening: {edit_file_path}")
        print(f"[Edit] Applying {len(edits)} operations...")
        result = edit_pptx(str(edit_file_path), edits, str(output_path), gemini)
        print(f"\n[Done] Edited copy saved to: {result}")
        return

    # ── Determine content source ─────────────────────────────────────────────
    if args.json_file:
        print(f"[Mode] Loading pre-generated JSON from: {args.json_file}")
        with open(args.json_file, "r", encoding="utf-8") as f:
            content = json.load(f)
    else:
        # Build topic string
        if args.prompt_file:
            with open(args.prompt_file, "r", encoding="utf-8") as f:
                topic = f.read()
        elif args.topic:
            topic = args.topic
        else:
            # Interactive mode
            print("No topic provided. Enter your presentation topic / instructions:")
            topic = input("> ").strip()
            if not topic:
                print("No topic given. Exiting.")
                sys.exit(1)

        print(f"\n[Gemini] Generating content for: {topic[:80]}...")
        gemini = GeminiClient(api_key=args.api_key)
        prompt = build_gemini_prompt(topic, context=args.context)

        # Save the prompt for debugging
        prompt_log = OUTPUT_DIR / "last_prompt.txt"
        prompt_log.write_text(prompt, encoding="utf-8")
        print(f"[Gemini] Prompt saved to: {prompt_log}")

        content = gemini.generate_json(prompt)

        # Save the generated JSON for reuse
        json_log = OUTPUT_DIR / "last_content.json"
        json_log.write_text(json.dumps(content, indent=2), encoding="utf-8")
        print(f"[Gemini] Content JSON saved to: {json_log}")

    # ── Build the PPT ────────────────────────────────────────────────────────
    template_path = Path(args.template)
    if not template_path.exists():
        print(f"[ERROR] Template not found at: {template_path}")
        print(f"Place the KaarTech Corporate PPT Templates file at:")
        print(f"  {template_path}")
        sys.exit(1)

    # Determine output path
    title = content.get("presentation_title", "KaarTech_Presentation")
    safe_title = "".join(c for c in title if c.isalnum() or c in " _-")[:40]
    safe_title = safe_title.strip().replace(" ", "_") or "presentation"

    if args.output:
        out_path = OUTPUT_DIR / args.output
    else:
        out_path = OUTPUT_DIR / f"{safe_title}.pptx"
        # Avoid overwriting — append suffix if file exists
        if out_path.exists():
            counter = 1
            while out_path.exists():
                out_path = OUTPUT_DIR / f"{safe_title}_{counter}.pptx"
                counter += 1

    result = build(
        prs_path=template_path,
        content=content,
        output_path=out_path,
        debug=args.debug,
    )

    print(f"\n[OK] Presentation generation complete!")
    print(f"     Output: {result}")


if __name__ == "__main__":
    main()
