"""
================================================================================
KAARTECH PPT EDITOR  —  ppt_editor.py
================================================================================
Author  : Thiso Vallabadass K (AID-902015)
Purpose : Edit any existing .pptx file from the command line.
          Works ALONGSIDE ppt_generator.py — does NOT touch generate logic.

COMMANDS:
  inspect   — show every slide, shape, and text in a .pptx
  edit      — apply one or more changes from a JSON instructions file
  chat      — interactive mode: describe what you want in plain English,
              Gemini translates it to edit operations and applies them

USAGE:
  python ppt_editor.py inspect  --file "THISO_L2_Portal.pptx"
  python ppt_editor.py edit     --file "THISO_L2_Portal.pptx" --edits edits.json --output "v2.pptx"
  python ppt_editor.py edit     --file "THISO_L2_Portal.pptx" --edits '[{"op":"title","slide":1,"text":"New Title"}]' --output "v2.pptx"
  python ppt_editor.py chat     --file "THISO_L2_Portal.pptx" --output "v2.pptx"

DEPENDENCIES:
  pip install python-pptx google-generativeai python-dotenv lxml
================================================================================
"""

import os
import sys
import json
import copy
import shutil
import argparse
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import lxml.etree as etree

# ── Load .env for GEMINI_API_KEY ──────────────────────────────────────────────
try:
    from dotenv import load_dotenv
    # Look for .env in same dir, parent dir, or cwd
    for _env in [Path(__file__).parent / ".env",
                 Path(__file__).parent.parent / ".env",
                 Path.cwd() / ".env"]:
        if _env.exists():
            load_dotenv(_env)
            break
except ImportError:
    pass

# ── Brand colors (kept consistent with generator) ─────────────────────────────
KAAR_RED   = "C0202B"
KAAR_DARK  = "1A1A2E"
KAAR_WHITE = "FFFFFF"
KAAR_GRAY  = "585858"

OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)


# ═════════════════════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ═════════════════════════════════════════════════════════════════════════════

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _xml_set_text(shape, text: str,
                  font_name: str = None, size_pt: float = None,
                  bold: bool = None, color_hex: str = None):
    """
    Replace all text in a shape with explicit XML-level font overrides.
    This is the only reliable way to set bold/color in python-pptx —
    setting run.font.bold alone is often ignored due to theme inheritance.
    Supports multi-line text: split on '\\n', each line becomes a paragraph.
    """
    if not shape.has_text_frame:
        return

    txBody = shape.text_frame._txBody
    # Remove all existing <a:p> elements
    for p in list(txBody.findall(f"{{{NS}}}p")):
        txBody.remove(p)

    lines = str(text).split("\n") if text else [""]
    for line in lines:
        p_el = etree.SubElement(txBody, f"{{{NS}}}p")
        r_el = etree.SubElement(p_el,  f"{{{NS}}}r")
        rPr  = etree.SubElement(r_el,  f"{{{NS}}}rPr",
                                 attrib={"lang": "en-US", "dirty": "0"})
        if bold is not None:
            rPr.set("b", "1" if bold else "0")
        if size_pt is not None:
            rPr.set("sz", str(int(size_pt * 100)))
        if font_name:
            etree.SubElement(rPr, f"{{{NS}}}latin",
                             attrib={"typeface": font_name})
        if color_hex:
            sf = etree.SubElement(rPr, f"{{{NS}}}solidFill")
            etree.SubElement(sf, f"{{{NS}}}srgbClr",
                             attrib={"val": color_hex.lstrip("#")})
        t_el = etree.SubElement(r_el, f"{{{NS}}}t")
        t_el.text = line


def _xml_append_paragraph(shape, text: str,
                           font_name: str = None, size_pt: float = None,
                           bold: bool = None, color_hex: str = None):
    """Append one paragraph to an existing text frame."""
    if not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    p_el = etree.SubElement(txBody, f"{{{NS}}}p")
    r_el = etree.SubElement(p_el,  f"{{{NS}}}r")
    rPr  = etree.SubElement(r_el,  f"{{{NS}}}rPr",
                             attrib={"lang": "en-US", "dirty": "0"})
    if bold is not None:
        rPr.set("b", "1" if bold else "0")
    if size_pt is not None:
        rPr.set("sz", str(int(size_pt * 100)))
    if font_name:
        etree.SubElement(rPr, f"{{{NS}}}latin",
                         attrib={"typeface": font_name})
    if color_hex:
        sf = etree.SubElement(rPr, f"{{{NS}}}solidFill")
        etree.SubElement(sf, f"{{{NS}}}srgbClr",
                         attrib={"val": color_hex.lstrip("#")})
    t_el = etree.SubElement(r_el, f"{{{NS}}}t")
    t_el.text = str(text)


def _find_shape(slide, name: str = None, ph_idx: int = None):
    """
    Find a shape by name (partial, case-insensitive) or placeholder index.
    Returns None if not found — never crashes.
    """
    if ph_idx is not None:
        for s in slide.placeholders:
            try:
                if s.placeholder_format.idx == ph_idx:
                    return s
            except Exception:
                pass
    if name:
        nl = name.lower()
        for s in slide.shapes:
            if nl in s.name.lower():
                return s
    return None


def _largest_text_shape(slide):
    """Return the text shape with the largest area on the slide."""
    best, best_area = None, 0
    for s in slide.shapes:
        if s.has_text_frame:
            area = s.width * s.height
            if area > best_area:
                best_area = area
                best = s
    return best


def _replace_in_runs(shape, find: str, replace: str) -> int:
    """Find-replace in all runs of a shape. Returns count of replacements."""
    if not shape.has_text_frame:
        return 0
    count = 0
    for para in shape.text_frame.paragraphs:
        if find in para.text:
            if not para.runs:
                para.text = para.text.replace(find, replace)
                count += 1
            else:
                replaced_in_runs = False
                for run in para.runs:
                    if find in run.text:
                        run.text = run.text.replace(find, replace)
                        count += 1
                        replaced_in_runs = True
                if not replaced_in_runs:
                    para.text = para.text.replace(find, replace)
                    count += 1
    return count


# ═════════════════════════════════════════════════════════════════════════════
#  INSPECT
# ═════════════════════════════════════════════════════════════════════════════

def inspect_pptx(file_path: str) -> dict:
    """
    Return full slide-by-slide structure of any .pptx as a dict.
    Printed to console and also returned for programmatic use.
    """
    prs = Presentation(file_path)
    result = {
        "file": Path(file_path).name,
        "total_slides": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "layout_name": slide.slide_layout.name,
            "shapes": []
        }
        for shape in slide.shapes:
            ph_idx = None
            if shape.is_placeholder:
                try:
                    ph_idx = shape.placeholder_format.idx
                except Exception:
                    pass
            full_text = shape.text_frame.text if shape.has_text_frame else ""
            font_info = {}
            if shape.has_text_frame:
                try:
                    para = shape.text_frame.paragraphs[0]
                    run  = para.runs[0] if para.runs else None
                    if run:
                        font_info = {
                            "name":  run.font.name,
                            "size":  round(run.font.size / 12700) if run.font.size else None,
                            "bold":  run.font.bold,
                            "color": str(run.font.color.rgb) if (run.font.color and run.font.color.type) else None,
                        }
                except Exception:
                    pass
            slide_info["shapes"].append({
                "name":         shape.name,
                "ph_idx":       ph_idx,
                "has_text":     shape.has_text_frame,
                "text_preview": full_text[:100].replace("\n", " \n "),
                "full_text":    full_text,
                "font":         font_info,
            })
        result["slides"].append(slide_info)

    return result


def print_inspect(result: dict):
    """Pretty-print the inspect result to the terminal."""
    print(f"\n{'='*70}")
    print(f"  FILE : {result['file']}")
    print(f"  SLIDES: {result['total_slides']}")
    print(f"{'='*70}")
    for s in result["slides"]:
        print(f"\n  Slide {s['slide_number']:>2}  [{s['layout_name']}]")
        for sh in s["shapes"]:
            if not sh["has_text"]:
                continue
            ph = f" ph={sh['ph_idx']}" if sh["ph_idx"] is not None else ""
            font = ""
            if sh["font"]:
                f = sh["font"]
                font = f"  font={f.get('name','?')} {f.get('size','?')}pt bold={f.get('bold','?')}"
            print(f"    * {sh['name']!r:35s}{ph}")
            if sh["text_preview"]:
                print(f"        text: {sh['text_preview'][:80]!r}{font}")


# ═════════════════════════════════════════════════════════════════════════════
#  EDIT OPERATIONS
# ═════════════════════════════════════════════════════════════════════════════

def apply_edits(prs: Presentation, edits: list) -> list:
    """
    Apply a list of edit operations to an open Presentation object in-place.
    Returns a list of result messages (one per operation).

    OPERATION REFERENCE
    -------------------
    Every operation is a dict with an "op" key. Valid values:

    "title"       — Change the title of a slide (placeholder idx=0)
    "section"     — Change a section divider slide text (largest shape)
    "text"        — Replace all text in a specific shape
    "replace"     — Find/replace text in one slide (preserves formatting)
    "replace_all" — Find/replace across ALL slides
    "append"      — Append a line/bullet to a shape
    "add_slide"   — Insert a new slide after a given position
    "delete_slide"— Remove a slide (processed last, reverse order)
    "font_slide"  — Change font/size/bold/color on ALL text in one slide
    """
    slides = prs.slides
    results = []

    def get_slide(n):
        idx = int(n) - 1
        if idx < 0 or idx >= len(slides):
            raise ValueError(f"Slide {n} out of range (total {len(slides)})")
        return slides[idx]

    def find_layout(name_fragment):
        nl = name_fragment.lower()
        for lay in prs.slide_layouts:
            if nl in lay.name.lower():
                return lay
        return prs.slide_layouts[0]

    # Split out delete ops — apply them last in reverse order
    delete_ops   = [e for e in edits if e.get("op") == "delete_slide"]
    regular_ops  = [e for e in edits if e.get("op") != "delete_slide"]

    for edit in regular_ops:
        op = edit.get("op", "").lower()
        try:
            # ── title ────────────────────────────────────────────────────────
            if op == "title":
                slide = get_slide(edit["slide"])
                shape = _find_shape(slide, ph_idx=0) or \
                        _find_shape(slide, name="Title")
                if not shape:
                    shape = _largest_text_shape(slide)
                if shape:
                    _xml_set_text(shape, edit["text"],
                                  font_name=edit.get("font", "Arial"),
                                  size_pt=edit.get("size", 28),
                                  bold=edit.get("bold", True),
                                  color_hex=edit.get("color", KAAR_RED))
                    results.append(f"[OK] title slide={edit['slide']}: {edit['text'][:50]!r}")
                else:
                    results.append(f"[WARN] title slide={edit['slide']}: no title shape found")

            # ── section ──────────────────────────────────────────────────────
            elif op == "section":
                slide = get_slide(edit["slide"])
                shape = _largest_text_shape(slide)
                if shape:
                    _xml_set_text(shape, edit["text"].upper(),
                                  font_name=edit.get("font", "Arial"),
                                  size_pt=edit.get("size", 36),
                                  bold=edit.get("bold", True),
                                  color_hex=edit.get("color", KAAR_WHITE))
                    results.append(f"[OK] section slide={edit['slide']}: {edit['text']!r}")
                else:
                    results.append(f"[WARN] section slide={edit['slide']}: no shape found")

            # ── text ─────────────────────────────────────────────────────────
            elif op == "text":
                slide = get_slide(edit["slide"])
                ph_idx = edit.get("ph")
                shape = _find_shape(slide,
                                    name=edit.get("shape"),
                                    ph_idx=int(ph_idx) if ph_idx is not None else None)
                if shape:
                    _xml_set_text(shape, edit["text"],
                                  font_name=edit.get("font"),
                                  size_pt=edit.get("size"),
                                  bold=edit.get("bold"),
                                  color_hex=edit.get("color"))
                    results.append(f"[OK] text slide={edit['slide']} shape={edit.get('shape') or f'ph={ph_idx}'}")
                else:
                    results.append(f"[WARN] text slide={edit['slide']}: shape {edit.get('shape')!r} not found")

            # ── replace ───────────────────────────────────────────────────────
            elif op == "replace":
                slide = get_slide(edit["slide"])
                total = 0
                for shape in slide.shapes:
                    total += _replace_in_runs(shape, edit["find"], edit["replace"])
                results.append(f"[OK] replace slide={edit['slide']}: {edit['find']!r} -> {edit['replace']!r} ({total} hits)")

            # ── replace_all ───────────────────────────────────────────────────
            elif op == "replace_all":
                total = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        total += _replace_in_runs(shape, edit["find"], edit["replace"])
                results.append(f"[OK] replace_all: {edit['find']!r} -> {edit['replace']!r} ({total} hits)")

            # ── append ────────────────────────────────────────────────────────
            elif op == "append":
                slide = get_slide(edit["slide"])
                ph_idx = edit.get("ph")
                shape = _find_shape(slide,
                                    name=edit.get("shape"),
                                    ph_idx=int(ph_idx) if ph_idx is not None else None)
                if shape:
                    _xml_append_paragraph(shape, edit["text"],
                                          font_name=edit.get("font", "Calibri"),
                                          size_pt=edit.get("size", 13),
                                          bold=edit.get("bold", False),
                                          color_hex=edit.get("color", KAAR_GRAY))
                    results.append(f"[OK] append slide={edit['slide']}")
                else:
                    results.append(f"[WARN] append slide={edit['slide']}: shape not found")

            # ── add_slide ─────────────────────────────────────────────────────
            elif op == "add_slide":
                layout = find_layout(edit.get("layout", "Content Slide_3"))
                new_slide = prs.slides.add_slide(layout)

                # Move to correct position
                xml_slides = prs.slides._sldIdLst
                last = xml_slides[-1]
                xml_slides.remove(last)
                after = int(edit.get("after", len(xml_slides)))
                insert_pos = min(after, len(xml_slides))
                xml_slides.insert(insert_pos, last)

                # Fill title (ph_idx=0)
                if edit.get("title"):
                    t_shape = _find_shape(new_slide, ph_idx=0) or \
                              _find_shape(new_slide, name="Title")
                    if t_shape:
                        _xml_set_text(t_shape, edit["title"],
                                      font_name=edit.get("font", "Arial"),
                                      size_pt=edit.get("size", 28),
                                      bold=True,
                                      color_hex=KAAR_RED)

                # Fill body (ph_idx=1)
                if edit.get("body"):
                    b_shape = _find_shape(new_slide, ph_idx=1) or \
                              _find_shape(new_slide, name="Content")
                    if b_shape:
                        _xml_set_text(b_shape, edit["body"],
                                      font_name="Calibri",
                                      size_pt=13,
                                      bold=False,
                                      color_hex=KAAR_GRAY)

                results.append(f"[OK] add_slide after={edit.get('after')} layout={layout.name!r}")

            # ── font_slide ────────────────────────────────────────────────────
            elif op == "font_slide":
                slide = get_slide(edit["slide"])
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    current_text = shape.text_frame.text
                    if current_text.strip():
                        _xml_set_text(shape, current_text,
                                      font_name=edit.get("font"),
                                      size_pt=edit.get("size"),
                                      bold=edit.get("bold"),
                                      color_hex=edit.get("color"))
                results.append(f"[OK] font_slide slide={edit['slide']}")

            else:
                results.append(f"[SKIP] unknown op: {op!r}")

        except Exception as e:
            results.append(f"[ERROR] op={op} slide={edit.get('slide','?')}: {e}")

    # ── Delete ops in reverse order ───────────────────────────────────────────
    for edit in sorted(delete_ops, key=lambda e: int(e.get("slide", 0)), reverse=True):
        try:
            idx = int(edit["slide"]) - 1
            xml_slides = prs.slides._sldIdLst
            if 0 <= idx < len(xml_slides):
                xml_slides.remove(xml_slides[idx])
                results.append(f"[OK] delete_slide slide={edit['slide']}")
            else:
                results.append(f"[WARN] delete_slide: slide {edit['slide']} out of range")
        except Exception as e:
            results.append(f"[ERROR] delete_slide slide={edit.get('slide','?')}: {e}")

    return results


# ═════════════════════════════════════════════════════════════════════════════
#  GEMINI CLIENT (minimal — only what the editor needs)
# ═════════════════════════════════════════════════════════════════════════════

class GeminiEditor:
    """
    Lightweight Gemini wrapper for the editor.
    Translates plain-English instructions -> edit operations JSON.
    """
    MODEL = "gemini-2.0-flash"

    def __init__(self, api_key: str = None):
        import google.generativeai as genai
        key = api_key or os.getenv("GEMINI_API_KEY")
        if not key:
            raise ValueError(
                "GEMINI_API_KEY not found. Set it in .env or pass --api-key."
            )
        genai.configure(api_key=key)
        self.model = genai.GenerativeModel(
            model_name=self.MODEL,
            generation_config=genai.GenerationConfig(
                response_mime_type="application/json",
                temperature=0.2,
                max_output_tokens=4096,
            )
        )

    def instructions_to_edits(self, user_instruction: str,
                               slide_structure: dict) -> list:
        """
        Given a plain-English instruction and the current slide structure,
        return a list of edit operations the editor can apply.
        """
        structure_summary = "\n".join(
            f"  Slide {s['slide_number']}: [{s['layout_name']}] "
            + ", ".join(
                f"{sh['name']}={sh['text_preview'][:40]!r}"
                for sh in s["shapes"] if sh["has_text"] and sh["text_preview"]
            )
            for s in slide_structure["slides"]
        )

        prompt = f"""
You are a PowerPoint editing assistant for Kaar Technologies.
The user wants to edit a presentation. Translate their instruction into
a JSON array of edit operations. Return ONLY the JSON array. No markdown.

CURRENT PRESENTATION STRUCTURE:
{structure_summary}

USER INSTRUCTION:
{user_instruction}

AVAILABLE OPERATIONS (choose the right ones):

{{"op":"title","slide":N,"text":"...","font":"Arial","size":28,"bold":true,"color":"C0202B"}}
  -> Change the title of slide N

{{"op":"section","slide":N,"text":"...","size":36,"color":"FFFFFF"}}
  -> Change a section divider slide text (Slide Splitter layout)

{{"op":"text","slide":N,"shape":"shape name (partial ok)","text":"...","font":"Calibri","size":13,"bold":false,"color":"585858"}}
  -> Replace full content of a specific shape

{{"op":"replace","slide":N,"find":"old text","replace":"new text"}}
  -> Find/replace in one slide (preserves formatting)

{{"op":"replace_all","find":"old text","replace":"new text"}}
  -> Find/replace across all slides (e.g. update dates, names)

{{"op":"append","slide":N,"shape":"shape name","text":"\\u2022 New bullet","font":"Calibri","size":13}}
  -> Add a line at the end of a shape

{{"op":"add_slide","after":N,"layout":"Content Slide_3 - Title and Content","title":"...","body":"..."}}
  -> Add a new slide after slide N

{{"op":"delete_slide","slide":N}}
  -> Delete slide N (processed last so other slide numbers stay valid)

RULES:
- Use "replace_all" for things like changing a date, name, or ID everywhere
- Use "title" for slide titles (placeholder idx=0)
- Use "section" for Slide Splitter / section divider slides
- Use "text" when you need to replace the full content of a content box
- Use "append" when you want to ADD something to existing content
- For shape names, use partial names from the structure above (e.g. "Content Placeholder")
- color values are 6-char hex WITHOUT the # sign
- Return a JSON array: [ {{op1}}, {{op2}}, ... ]
- If the instruction is ambiguous, make a reasonable interpretation
"""

        for attempt in range(3):
            try:
                response = self.model.generate_content(prompt)
                raw = response.text.strip()
                # Strip markdown fences if any
                if "```" in raw:
                    for part in raw.split("```"):
                        cleaned = part.strip().lstrip("json").strip()
                        if cleaned.startswith("["):
                            raw = cleaned
                            break
                start = raw.find("[")
                end   = raw.rfind("]") + 1
                if start != -1 and end > start:
                    raw = raw[start:end]
                ops = json.loads(raw)
                if isinstance(ops, list):
                    return ops
            except Exception as e:
                if attempt == 2:
                    raise ValueError(f"Gemini failed to produce edit ops: {e}")
                prompt = "CRITICAL: last response was not a valid JSON array. " + prompt

        return []


# ═════════════════════════════════════════════════════════════════════════════
#  CORE WORKFLOW FUNCTIONS
# ═════════════════════════════════════════════════════════════════════════════

def resolve_path(file_arg: str) -> Path:
    """
    Resolve a file argument to an absolute Path.
    Looks in: exact path -> output/ -> cwd -> script dir parent/output
    """
    p = Path(file_arg)
    if p.is_absolute() and p.exists():
        return p
    # Try relative to cwd
    if (Path.cwd() / file_arg).exists():
        return Path.cwd() / file_arg
    # Try output dir
    if (OUTPUT_DIR / file_arg).exists():
        return OUTPUT_DIR / file_arg
    # Try as-is
    if p.exists():
        return p.resolve()
    raise FileNotFoundError(
        f"File not found: {file_arg}\n"
        f"Looked in: cwd={Path.cwd()}, output={OUTPUT_DIR}"
    )


def run_inspect(args):
    """python ppt_editor.py inspect --file FILE [--json]"""
    file_path = resolve_path(args.file)
    print(f"\n[Inspect] Reading: {file_path}")
    result = inspect_pptx(str(file_path))
    print_inspect(result)

    # Save JSON structure
    json_out = OUTPUT_DIR / "inspected_structure.json"
    json_out.write_text(json.dumps(result, indent=2, ensure_ascii=False),
                        encoding="utf-8")
    print(f"\n[Inspect] Full structure saved to: {json_out}")

    if getattr(args, "json", False):
        print(json.dumps(result, indent=2))


def run_edit(args):
    """python ppt_editor.py edit --file FILE --edits EDITS [--output OUT]"""
    file_path = resolve_path(args.file)

    # Parse edits: JSON string or path to .json file
    edits_input = args.edits.strip()
    # Only try Path.exists() if it looks like a filename (short, doesn't start with [)
    _looks_like_path = len(edits_input) < 260 and not edits_input.strip().startswith("[")
    if _looks_like_path:
        try:
            _p = Path(edits_input)
            if _p.exists():
                with open(_p, "r", encoding="utf-8-sig") as f:
                    edits = json.load(f)
            else:
                edits = json.loads(edits_input)
        except OSError:
            edits = json.loads(edits_input)
    else:
        edits = json.loads(edits_input)

    if not isinstance(edits, list):
        # Allow a single op dict without array
        edits = [edits]

    # Determine output path
    out_name = args.output or f"edited_{file_path.stem}.pptx"
    if not out_name.endswith(".pptx"):
        out_name += ".pptx"
    # If output has no directory component, put it in output/
    out_path = Path(out_name)
    if not out_path.is_absolute() and "/" not in out_name and "\\" not in out_name:
        out_path = OUTPUT_DIR / out_name

    # Copy original -> output (never modify original)
    shutil.copy2(str(file_path), str(out_path))
    print(f"\n[Edit] Source  : {file_path}")
    print(f"[Edit] Output  : {out_path}")
    print(f"[Edit] Applying {len(edits)} operation(s)...\n")

    prs = Presentation(str(out_path))
    results = apply_edits(prs, edits)
    prs.save(str(out_path))

    for r in results:
        print(f"  {r}")

    size_kb = out_path.stat().st_size // 1024
    print(f"\n[Done] Saved: {out_path}  ({size_kb} KB)")


def run_chat(args):
    """
    python ppt_editor.py chat --file FILE [--output OUT] [--api-key KEY]
    Interactive loop: type instructions in plain English, Gemini converts
    them to edit ops and applies them. Type 'done' to save and exit.
    """
    file_path = resolve_path(args.file)
    out_name  = args.output or f"edited_{file_path.stem}.pptx"
    if not out_name.endswith(".pptx"):
        out_name += ".pptx"
    out_path = Path(out_name)
    if not out_path.is_absolute() and "/" not in out_name and "\\" not in out_name:
        out_path = OUTPUT_DIR / out_name

    # Copy original to working copy
    shutil.copy2(str(file_path), str(out_path))
    print(f"\n[Chat] Working copy: {out_path}")
    print("[Chat] Type your edit instructions in plain English.")
    print("[Chat] Type 'inspect' to see current slide structure.")
    print("[Chat] Type 'undo' to reload the last saved version.")
    print("[Chat] Type 'done' or 'save' to finish and exit.\n")

    gemini = GeminiEditor(api_key=getattr(args, "api_key", None))

    while True:
        try:
            instruction = input("What would you like to change? > ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\n[Chat] Interrupted.")
            break

        if not instruction:
            continue

        if instruction.lower() in ("done", "save", "exit", "quit"):
            print(f"\n[Chat] Saved: {out_path}")
            break

        if instruction.lower() == "inspect":
            result = inspect_pptx(str(out_path))
            print_inspect(result)
            continue

        if instruction.lower() == "undo":
            shutil.copy2(str(file_path), str(out_path))
            print("[Chat] Reverted to original.")
            continue

        # Ask Gemini to translate instruction -> ops
        print("[Chat] Thinking...")
        try:
            structure = inspect_pptx(str(out_path))
            ops = gemini.instructions_to_edits(instruction, structure)
        except Exception as e:
            print(f"[Chat] Gemini error: {e}")
            continue

        if not ops:
            print("[Chat] Could not generate edit operations. Try rephrasing.")
            continue

        print(f"[Chat] Generated {len(ops)} operation(s):")
        for op in ops:
            print(f"  {json.dumps(op)}")

        confirm = input("Apply? [Y/n] > ").strip().lower()
        if confirm in ("n", "no"):
            print("[Chat] Skipped.")
            continue

        prs = Presentation(str(out_path))
        results = apply_edits(prs, ops)
        prs.save(str(out_path))
        for r in results:
            print(f"  {r}")
        print(f"[Chat] Saved. ({out_path.stat().st_size // 1024} KB)")


# ═════════════════════════════════════════════════════════════════════════════
#  CLI
# ═════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        prog="ppt_editor",
        description="KaarTech PPT Editor \u2014 edit any .pptx from the command line",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
EXAMPLES:

  # See what's in a presentation
  python ppt_editor.py inspect --file THISO_L2_Portal.pptx

  # Update slide 1 title and replace date everywhere
  python ppt_editor.py edit --file THISO_L2_Portal.pptx --output v2.pptx ^
    --edits "[{\\"op\\":\\"title\\",\\"slide\\":1,\\"text\\":\\"PORTAL REVIEW-L3\\"},{\\"op\\":\\"replace_all\\",\\"find\\":\\"07-06-2026\\",\\"replace\\":\\"12-06-2026\\"}]"

  # Load edits from a JSON file
  python ppt_editor.py edit --file THISO_L2_Portal.pptx --edits my_edits.json --output v2.pptx

  # Interactive AI chat mode
  python ppt_editor.py chat --file THISO_L2_Portal.pptx --output v2.pptx
"""
    )

    sub = parser.add_subparsers(dest="command", required=True)

    # ── inspect ───────────────────────────────────────────────────────────────
    p_inspect = sub.add_parser("inspect",
        help="Show all slides, shapes, and text in a .pptx file")
    p_inspect.add_argument("--file", "-f", required=True,
        help="Path or filename of the .pptx to inspect")
    p_inspect.add_argument("--json", action="store_true",
        help="Also print raw JSON structure to stdout")

    # ── edit ──────────────────────────────────────────────────────────────────
    p_edit = sub.add_parser("edit",
        help="Apply edit operations to a .pptx file")
    p_edit.add_argument("--file", "-f", required=True,
        help="Source .pptx file to edit")
    p_edit.add_argument("--edits", "-e", required=True,
        help="JSON array of edit operations (string or path to .json file)")
    p_edit.add_argument("--output", "-o",
        help="Output filename (saved to output/ if no path given)")

    # ── chat ──────────────────────────────────────────────────────────────────
    p_chat = sub.add_parser("chat",
        help="Interactive AI-powered editing session")
    p_chat.add_argument("--file", "-f", required=True,
        help="Source .pptx file to edit")
    p_chat.add_argument("--output", "-o",
        help="Output filename for the edited copy")
    p_chat.add_argument("--api-key", "-k",
        help="Gemini API key (overrides .env)")

    args = parser.parse_args()

    if args.command == "inspect":
        run_inspect(args)
    elif args.command == "edit":
        run_edit(args)
    elif args.command == "chat":
        run_chat(args)


if __name__ == "__main__":
    main()
